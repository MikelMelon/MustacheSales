# generate_modules_4_8.py
# Script to generate Module 4-8 PPTX files using an existing PPTX as template.
# Usage:
# 1) pip install python-pptx
# 2) Place this script next to your template PPTX (e.g. "1 МОДУЛЬ Введение в музыкальный продакшен.pptx" or "Ableton_Live.pptx")
# 3) python generate_modules_4_8.py --template "path/to/template.pptx" --out generated_ppts

import argparse
import os
from pptx import Presentation
from pptx.util import Pt

# --- Настройка контента модулей (можно править) ---
modules = {
    4: {
        "title": "Модуль 4: Аранжировка",
        "slides": [
            {"title": "Введение в аранжировку", "bullets": [
                "Основные элементы трека (интро, куплет, припев, бридж, аутро)",
                "Роль динамики и развития",
                "Примеры классических структур"]},
            {"title": "Работа с идеями", "bullets": [
                "Как превращать лупы и мелодии в части трека",
                "Через композицию — контраст и повтор",
                "Построение музыкального нарратива"]},
            {"title": "Построение переходов", "bullets": [
                "Преходы: фильтры, ударные, падсы",
                "Автоматизация параметров для плавных переходов",
                "Использование эффектов для акцентов"]},
            {"title": "Структурирование драматургии", "bullets": [
                "Кульминация и релаксация — где и почему",
                "Размещение вокала и основных тем",
                "Микс баланса для разных частей"]},
            {"title": "Работа с энергией трека", "bullets": [
                "Использование частот и динамики для роста/спада",
                "Гайдлайны: когда добавлять/убирать элементы",
                "Техники удержания интереса слушателя"]},
            {"title": "Практическое задание", "bullets": [
                "Развернуть 8-тактовый луп в полную структуру",
                "Создать 3 перехода между частями",
                "Подготовить короткий демо-аранж (1.5–2 мин)"]},
        ]
    },
    5: {
        "title": "Модуль 5: Саунд‑дизайн и эффекты",
        "slides": [
            {"title": "Основы звукового дизайна", "bullets": [
                "Типы синтеза: субтрактив, FM, wavetable, гранулярный",
                "Ключевые параметры: осцилляторы, фильтры, огибающие",
                "Рабочие примеры в Ableton (Wavetable, Analog)"]},
            {"title": "Создание баса и лидов", "bullets": [
                "Тональность и формантная обработка",
                "Слои и синхронизация фаз",
                "Приёмы для мощного баса"]},
            {"title": "Текстуры, падсы и атмосферные звуки", "bullets": [
                "Генерация богатых слоёв с помощью эффектов",
                "Реверберация и модуляция для глубины",
                "Использование Sampler/Simpler для текстур"]},
            {"title": "Эффекты — принципы и цепочки", "bullets": [
                "Порядок эффектов и влияние на звук",
                "Динамические эффекты: компрессия, transient shaper",
                "Модуляционные эффекты: chorus, phaser, flanger"]},
            {"title": "Creative FX workflows", "bullets": [
                "Freeze & resample — как создать уникальную текстуру",
                "Sidechain и ducking для движения",
                "Автоматизация эффектов для живости звучания"]},
            {"title": "Практическое задание", "bullets": [
                "Создать 3 уникальных пресета (бас, пад, эффектный лид)",
                "Показать цепочку эффектов и их параметры"]},
        ]
    },
    6: {
        "title": "Модуль 6: Сведение (Mixing)",
        "slides": [
            {"title": "Введение в сведение", "bullets": [
                "Цели свода: ясность, баланс, пространство",
                "Подготовка сессии: метки, группы, цветовая кодировка",
                "Сигнальные цепи и уровни"]},
            {"title": "Эквализация в контексте", "bullets": [
                "Удаление конфликтов частот (какой инструмент где живёт)",
                "Мастеринг эквалайзеры vs каналовые",
                "Практические примеры: бас и ударные"]},
            {"title": "Динамика: компрессия и сатурация", "bullets": [
                "Природа компрессии и основные настройки",
                "Параллельная компрессия и glue-компрессия",
                "Сатурация для гармонической насыщенности"]},
            {"title": "Панорамирование и глубина", "bullets": [
                "Создание ширины и места для элементов",
                "Reverb/Delay как инструмент глубины",
                "Автоматизация панорамы для движения"]},
            {"title": "Работа с группами и посылами", "bullets": [
                "Шины посыла: реверб, делэй, эффектные посылы",
                "Субмиксинг барабанов и инструментов",
                "Управление громкостью и бёдом"]},
            {"title": "Контроль громкости и метрики", "bullets": [
                "LUFS, True Peak — ориентиры для стриминга",
                "Использование анализаторов и референсов",
                "Практическое сравнение с референс-треком"]},
            {"title": "Практическое задание", "bullets": [
                "Свести короткий 2-мин трек: баланс, EQ, компрессия",
                "Подготовить 2 версии: для стрима и для подкаста"]},
        ]
    },
    7: {
        "title": "Модуль 7: Мастеринг",
        "slides": [
            {"title": "Что такое мастеринг", "bullets": [
                "Отличие мастеринга от сведения",
                "Задачи мастеринга: громкость, прозрачность, совместимость",
                "Последовательность мастеринга"]},
            {"title": "Подготовка к мастерингу", "bullets": [
                "Экспорт промежуточных файлов: стемы vs стерео",
                "Чистка пиков и выравнивание громкости",
                "Рекомендации по headroom (‑6…‑3 dB)"]},
            {"title": "Инструменты мастеринга", "bullets": [
                "Линейная и многополосная компрессия",
                "Лимитеры и True Peak контроль",
                "Стерео-ширина и коррекция частот"]},
            {"title": "Рабочий процесс мастеринга", "bullets": [
                "Референсные треки и конечные цели",
                "Чек-лист мастера: спектр, динамика, фазовые проблемы",
                "AB тестирование и экспорт форматов"]},
            {"title": "Проблемы и решения", "bullets": [
                "Частые ошибки (перепресс, потеря динамики)",
                "Как исправить узкий/грязный микс",
                "Советы по мастерингу в домашних условиях"]},
            {"title": "Практическое задание", "bullets": [
                "Замастерить 1-минутный фрагмент с референсом",
                "Подготовить финальную волну (WAV) и MP3"]},
        ]
    },
    8: {
        "title": "Модуль 8: Дистрибуция и продвижение",
        "slides": [
            {"title": "Каналы дистрибуции", "bullets": [
                "Платформы: стриминговые сервисы, магазины, Bandcamp",
                "Агрегаторы и их роль",
                "Выбор релизного формата"]},
            {"title": "Метаданные и права", "bullets": [
                "ISRC, титры, данные по авторам",
                "Авторские права и регистрация трека",
                "Музыкальные публикации и роялти"]},
            {"title": "Продвижение релиза", "bullets": [
                "Планирование релиза: пресейвы, тизеры, плейлисты",
                "Работа с прессой и блогерами",
                "Социальные сети: стратегия и контент-план"]},
            {"title": "Плейлисты и сетевое продвижение", "bullets": [
                "Как попасть в плейлисты: тактика и реальность",
                "Коллаборации и кросс-промо",
                "Реклама и таргет vs органика"]},
            {"title": "Монетизация и дальнейшие шаги", "bullets": [
                "Концерты, мерч, лайцензирование",
                "Создание карьеры: бренд и портфолио",
                "Аналитика: метрики успеха"]},
            {"title": "Практическое задание", "bullets": [
                "Составить план релиза для готового трека",
                "Подготовить метаданные и промо-план"]},
        ]
    }
}

# --- Конструктор презентаций ---
def remove_all_slides(prs):
    # Удаляет все слайды в презентации (manipulates the internal slide id list)
    slide_ids = [sld.id for sld in prs.slides._sldIdLst]
    # Can't manipulate list directly via public API; use xml
    xml_slides = prs.slides._sldIdLst
    for sldId in list(xml_slides):
        xml_slides.remove(sldId)

def add_title_and_bullets(prs, layout_index, title_text, bullets, notes_text=None):
    layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(layout)
    # Title
    try:
        title = slide.shapes.title
        title.text = title_text
    except Exception:
        # find first placeholder
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.type == 1:
                shape.text = title_text
                break
    # Content / bullets
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.type == 2:
            tf = shape.text_frame
            tf.clear()
            for i, b in enumerate(bullets):
                p = tf.add_paragraph()
                p.text = b
                p.level = 0
                p.font.size = Pt(18)
            break
    # Notes
    if notes_text:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes_text
    return slide

def create_module_pptx(module_number, module_data, template_path, out_dir):
    title = module_data["title"]
    out_name = f"Module_{module_number}_{title}.pptx".replace(" ", "_")
    out_path = os.path.join(out_dir, out_name)
    # Используем шаблон как basis
    prs = Presentation(template_path)
    # Очистим начальные слайды, чтобы добавить наши
    remove_all_slides(prs)
    # Добавляем титульный слайд (используем layout 0 если есть)
    # Попробуем layout 0 -> титул, layout 1 -> заголовок+контент
    for i, s in enumerate(module_data["slides"]):
        if i == 0:
            layout_idx = 0 if len(prs.slide_layouts) > 0 else 1
        else:
            layout_idx = 1 if len(prs.slide_layouts) > 1 else 0
        add_title_and_bullets(prs, layout_idx, s["title"], s.get("bullets", []), notes_text=s.get("notes", ""))
    prs.save(out_path)
    return out_path

def main():
    parser = argparse.ArgumentParser(description="Generate Module PPTX files from a template.")
    parser.add_argument("--template", "-t", required=True, help="Path to template PPTX (your module 1 PPTX)")
    parser.add_argument("--out", "-o", default="output_ppts", help="Output folder for generated PPTX files")
    args = parser.parse_args()

    template_path = args.template
    out_dir = args.out
    if not os.path.isfile(template_path):
        print("Ошибка: шаблон не найден:", template_path)
        return
    os.makedirs(out_dir, exist_ok=True)

    for num, data in modules.items():
        print("Генерирую:", num, data["title"])
        path = create_module_pptx(num, data, template_path, out_dir)
        print("Сохранено:", path)

    print("Готово. Проверьте папку:", out_dir)

if __name__ == "__main__":
    main()
